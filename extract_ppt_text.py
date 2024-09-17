import os
from pptx import Presentation
import win32com.client as win32

# Folder containing the PowerPoint files
ppt_folder = "path_to_your_folder"  # Update this with the path to your folder

# Function to convert .ppt to .pptx
def convert_ppt_to_pptx(ppt_file):
    powerpoint = win32.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True  # Keep PowerPoint visible since hiding is not allowed
    
    # Open the presentation
    ppt = powerpoint.Presentations.Open(ppt_file)
    
    # Convert .ppt to .pptx
    pptx_file = ppt_file.replace(".ppt", ".pptx")
    ppt.SaveAs(pptx_file, 24)  # 24 is the format ID for .pptx
    ppt.Close()  # Close the presentation
    
    # Quit PowerPoint to avoid keeping it open
    powerpoint.Quit()
    
    return pptx_file

# Iterate through all files in the folder
for file_name in os.listdir(ppt_folder):
    ppt_file = os.path.join(ppt_folder, file_name)

    # Convert .ppt to .pptx if necessary
    if file_name.endswith(".ppt"):
        ppt_file = convert_ppt_to_pptx(ppt_file)
        print(f"Converted {file_name} to {ppt_file}")
    
    if ppt_file.endswith(".pptx"):  # Process only .pptx files
        print(f"\nProcessing file: {file_name}")

        # Load the PowerPoint presentation
        prs = Presentation(ppt_file)

        # Iterate through slides to extract titles and text from shapes
        for slide_num, slide in enumerate(prs.slides):
            print(f"\nSlide {slide_num + 1}:")
            
            # Extract the slide's title (header)
            if slide.shapes.title:
                title = slide.shapes.title.text
                print(f"Title: {title}")
            else:
                print("Title: None")

            # Extract text from other shapes
            for shape_num, shape in enumerate(slide.shapes):
                if shape == slide.shapes.title:
                    continue  # Skip the title shape

                if not shape.has_text_frame:
                    continue  # Skip shapes without text

                text = "\n".join([paragraph.text for paragraph in shape.text_frame.paragraphs])
                print(f"Shape {shape_num + 1}: {text}")